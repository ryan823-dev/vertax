from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable

from docx import Document
from docx.enum.section import WD_SECTION
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor as DocxRGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt


ROOT = Path(__file__).resolve().parent.parent
SOURCE_MD = ROOT / "docs" / "RADAR_CURRENT_USAGE_GUIDE.md"
WORD_OUTPUT = ROOT / "docs" / "获客雷达当前版本使用说明.docx"
PPT_OUTPUT = ROOT / "docs" / "获客雷达当前版本使用说明_汇报版.pptx"

FONT_NAME = "Microsoft YaHei"
TITLE_COLOR = (11, 27, 43)
GOLD = (212, 175, 55)
LIGHT_BG = (252, 249, 242)
MUTED = (94, 106, 120)
DARK_BG = (11, 18, 32)
WHITE = (255, 255, 255)


@dataclass
class Block:
    kind: str
    value: object


@dataclass
class SectionNode:
    level: int
    title: str
    blocks: list[Block] = field(default_factory=list)
    children: list["SectionNode"] = field(default_factory=list)


def clean_inline_markdown(text: str) -> str:
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    return text.strip()


def parse_markdown_sections(text: str) -> SectionNode:
    root = SectionNode(level=0, title="ROOT")
    stack = [root]
    lines = text.splitlines()
    i = 0

    while i < len(lines):
        line = lines[i].rstrip()

        heading_match = re.match(r"^(#{1,6})\s+(.*)$", line)
        if heading_match:
            level = len(heading_match.group(1))
            title = clean_inline_markdown(heading_match.group(2))
            node = SectionNode(level=level, title=title)
            while stack and stack[-1].level >= level:
                stack.pop()
            stack[-1].children.append(node)
            stack.append(node)
            i += 1
            continue

        if not line.strip():
            i += 1
            continue

        if line.startswith("|"):
            table_lines = [line]
            j = i + 1
            while j < len(lines) and lines[j].startswith("|"):
                table_lines.append(lines[j].rstrip())
                j += 1
            rows = []
            for row_line in table_lines:
                row = [clean_inline_markdown(cell) for cell in row_line.strip("|").split("|")]
                rows.append(row)
            if len(rows) >= 2 and all(re.fullmatch(r"[:\- ]+", cell) for cell in rows[1]):
                rows.pop(1)
            stack[-1].blocks.append(Block("table", rows))
            i = j
            continue

        if re.match(r"^[-*]\s+", line):
            items = []
            j = i
            while j < len(lines):
                current = lines[j].rstrip()
                if re.match(r"^[-*]\s+", current):
                    items.append(clean_inline_markdown(re.sub(r"^[-*]\s+", "", current)))
                    j += 1
                    continue
                if not current.strip():
                    break
                if items:
                    items[-1] = f"{items[-1]} {clean_inline_markdown(current)}"
                    j += 1
                    continue
                break
            stack[-1].blocks.append(Block("ul", items))
            i = j
            continue

        if re.match(r"^\d+\.\s+", line):
            items = []
            j = i
            while j < len(lines):
                current = lines[j].rstrip()
                if re.match(r"^\d+\.\s+", current):
                    items.append(clean_inline_markdown(re.sub(r"^\d+\.\s+", "", current)))
                    j += 1
                    continue
                if not current.strip():
                    break
                if items:
                    items[-1] = f"{items[-1]} {clean_inline_markdown(current)}"
                    j += 1
                    continue
                break
            stack[-1].blocks.append(Block("ol", items))
            i = j
            continue

        paragraph_lines = [clean_inline_markdown(line)]
        j = i + 1
        while j < len(lines):
            current = lines[j].rstrip()
            if (
                not current.strip()
                or current.startswith("|")
                or re.match(r"^(#{1,6})\s+", current)
                or re.match(r"^[-*]\s+", current)
                or re.match(r"^\d+\.\s+", current)
            ):
                break
            paragraph_lines.append(clean_inline_markdown(current))
            j += 1
        paragraph = " ".join(part for part in paragraph_lines if part).strip()
        if paragraph:
            stack[-1].blocks.append(Block("paragraph", paragraph))
        i = j

    return root


def iter_sections(node: SectionNode) -> Iterable[SectionNode]:
    for child in node.children:
        yield child
        yield from iter_sections(child)


def find_section(root: SectionNode, title_startswith: str) -> SectionNode | None:
    for node in iter_sections(root):
        if node.title.startswith(title_startswith):
            return node
    return None


def set_run_font(run, size: int | None = None, bold: bool | None = None, color: tuple[int, int, int] | None = None) -> None:
    run.font.name = FONT_NAME
    run._element.rPr.rFonts.set(qn("w:eastAsia"), FONT_NAME)
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = DocxRGBColor(*color)


def format_doc_styles(doc: Document) -> None:
    normal = doc.styles["Normal"]
    normal.font.name = FONT_NAME
    normal._element.rPr.rFonts.set(qn("w:eastAsia"), FONT_NAME)
    normal.font.size = Pt(11)

    for style_name in ("Title", "Heading 1", "Heading 2", "Heading 3", "Heading 4", "List Bullet", "List Number"):
        style = doc.styles[style_name]
        style.font.name = FONT_NAME
        style._element.rPr.rFonts.set(qn("w:eastAsia"), FONT_NAME)


def add_doc_paragraph(doc: Document, text: str, style: str | None = None) -> None:
    p = doc.add_paragraph(style=style)
    run = p.add_run(text)
    set_run_font(run)


def export_word(root: SectionNode) -> None:
    doc = Document()
    format_doc_styles(doc)

    section = doc.sections[0]
    section.top_margin = Inches(0.75)
    section.bottom_margin = Inches(0.75)
    section.left_margin = Inches(0.9)
    section.right_margin = Inches(0.9)

    title = doc.add_paragraph(style="Title")
    title.alignment = 1
    run = title.add_run("获客雷达当前版本使用说明")
    set_run_font(run, size=22, bold=True, color=TITLE_COLOR)

    subtitle = doc.add_paragraph()
    subtitle.alignment = 1
    subtitle_run = subtitle.add_run("基于当前版本代码实现整理，可直接用于内部培训、客户交付和汇报。")
    set_run_font(subtitle_run, size=10, color=MUTED)

    doc.add_paragraph()

    for node in root.children:
        _write_section_to_doc(doc, node)

    doc.save(WORD_OUTPUT)


def _write_section_to_doc(doc: Document, node: SectionNode) -> None:
    heading_level = max(1, min(4, node.level - 1))
    heading = doc.add_heading(level=heading_level)
    run = heading.add_run(node.title)
    set_run_font(run, size=max(12, 17 - heading_level), bold=True, color=TITLE_COLOR)

    for block in node.blocks:
        if block.kind == "paragraph":
            add_doc_paragraph(doc, str(block.value))
        elif block.kind == "ul":
            for item in block.value:
                add_doc_paragraph(doc, str(item), style="List Bullet")
        elif block.kind == "ol":
            for item in block.value:
                add_doc_paragraph(doc, str(item), style="List Number")
        elif block.kind == "table":
            rows = block.value
            if not rows:
                continue
            table = doc.add_table(rows=len(rows), cols=len(rows[0]))
            table.style = "Table Grid"
            for row_idx, row in enumerate(rows):
                for col_idx, cell_text in enumerate(row):
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(cell_text)
                    for paragraph in cell.paragraphs:
                        for run in paragraph.runs:
                            set_run_font(run, size=10, bold=row_idx == 0)
            doc.add_paragraph()

    for child in node.children:
        _write_section_to_doc(doc, child)


def build_slide_plan(root: SectionNode) -> list[dict[str, object]]:
    sec1 = find_section(root, "1. 这套系统现在是怎么工作的")
    sec2 = find_section(root, "2. 使用前准备")
    sec41 = find_section(root, "4.1 总览页")
    sec42 = find_section(root, "4.2 目标客户画像")
    sec43 = find_section(root, "4.3 自动搜索")
    sec44 = find_section(root, "4.4 候选池")
    sec45 = find_section(root, "4.5 线索库")
    sec46 = find_section(root, "4.6 今日外联")
    sec47 = find_section(root, "4.7 采购机会")
    sec5 = find_section(root, "5. 当前版本最值得团队记住的几个规则")
    sec6 = find_section(root, "6. 推荐给团队的使用分工")
    sec7 = find_section(root, "7. 当前版本的边界与注意事项")
    sec8 = find_section(root, "8. 页面速查")
    sec9 = find_section(root, "9. 一句话总结")

    return [
        {
            "type": "title",
            "title": "获客雷达当前版本使用说明",
            "subtitle": "Word / PPT 导出版\n基于当前代码实现整理",
        },
        {
            "title": "系统工作方式",
            "bullets": first_list_items(sec1, 6),
        },
        {
            "title": "使用前准备与标准路径",
            "bullets": [
                *first_list_items(sec2, 4),
                "第一次启动：同步画像 -> 跑搜索 -> 审核候选 -> 导入线索 -> 外联推进",
                "日常运营：先看今日外联，再消化候选池，最后在线索库继续推进",
            ],
        },
        {
            "title": "总览页",
            "subtitle": "/customer/radar",
            "bullets": [
                "用于查看雷达流水线状态、统计卡片和各模块快捷入口",
                "智能获客助手支持自然语言描述目标客户",
                "当前实现会创建一份持续搜索配置 RadarSearchProfile",
                "默认启用 MAPS 与 DIRECTORY，并默认开启 autoQualify / autoEnrich",
            ],
        },
        {
            "title": "目标客户画像",
            "subtitle": "/customer/radar/targeting",
            "bullets": [
                "维护 TargetingSpec：系统按什么客户特征找客户",
                "维护 ChannelMap：系统优先从哪些渠道、按什么思路去找",
                "支持从知识引擎同步最新画像与渠道图谱",
                "支持手工编辑、版本历史、待确认问题查看",
                "建议重点检查行业、地区、关键词、排除规则和渠道优先级",
            ],
        },
        {
            "title": "自动搜索",
            "subtitle": "/customer/radar/search",
            "bullets": [
                "开始自动搜索：按当前画像立即执行一轮任务",
                "按最新画像重新搜索：以最新画像重建查询并执行",
                "暂停 / 继续自动执行：管理已有 RadarSearchProfile",
                "重要：持续自动跑需要先由用户启动一次，之后系统按调度规则自动执行，不需要每天重复点击",
                "页面重点用于看画像摘要、能力范围、数据源健康度和搜索配置状态",
                "要点：即时搜索和总览页的持续搜索配置不是同一动作",
            ],
        },
        {
            "title": "候选池",
            "subtitle": "/customer/radar/candidates",
            "bullets": [
                "负责审核系统发现但尚未进入正式线索库的对象",
                "支持按状态、层级和关键词筛选，并支持批量操作",
                "可执行分级、排除、批量丰富情报、导入线索库、导入采购机会",
                "推荐审核策略：先处理 NEW，再按 A / B / C / 排除分层",
                "候选池以审核与导入为主，持续外联尽量转到线索库处理",
            ],
        },
        {
            "title": "线索库",
            "subtitle": "/customer/radar/prospects",
            "bullets": [
                "正式推进阶段主战场：只有导入后的公司才进入正式跟进池",
                "支持联系人管理、公司富化、背调 dossier、外联草稿与发送",
                "支持保存外联包草稿 / 模板，并可保存到营销内容库",
                "可查看外联历史、营销内容建议，并支持导出 CSV",
                "建议动作：先补联系人，再生成外联包，发送后及时记录结果",
            ],
        },
        {
            "title": "今日外联与采购机会",
            "bullets": [
                "今日外联 /customer/radar/daily：把当天最值得联系的对象分为电话优先、邮件优先、待补全",
                "页面支持直接拨号、直接发邮、打开公司工作台，是日执行工作台",
                "采购机会 /customer/radar/opportunities：单独管理招投标和项目型机会",
                "机会支持按阶段过滤，并在 IDENTIFIED 到 WON / LOST 之间推进",
            ],
        },
        {
            "title": "使用规则",
            "bullets": collect_rule_points(sec5),
        },
        {
            "title": "推荐分工",
            "bullets": [
                "销售负责人：重点看今日外联、线索库、采购机会",
                "运营 / 增长负责人：重点看目标客户画像、自动搜索、候选池",
                "内容 / 营销负责人：重点看线索库中的营销内容建议与保存到内容库的资产",
            ],
        },
        {
            "title": "边界与注意事项",
            "bullets": first_list_items(sec7, 8),
        },
        {
            "title": "页面速查与总结",
            "bullets": [
                *table_to_points(sec8, 5),
                first_paragraph(sec9) or "一句话总结：把获客雷达当成团队的日常获客操作台，而不只是搜索框。",
            ],
        },
    ]


def first_list_items(section: SectionNode | None, limit: int) -> list[str]:
    if not section:
        return []
    items: list[str] = []
    for block in section.blocks:
        if block.kind in {"ul", "ol"}:
            for item in block.value:
                if len(items) >= limit:
                    return items
                items.append(str(item))
    return items


def first_paragraph(section: SectionNode | None) -> str | None:
    if not section:
        return None
    for block in section.blocks:
        if block.kind == "paragraph":
            return str(block.value)
    return None


def table_to_points(section: SectionNode | None, limit: int) -> list[str]:
    if not section:
        return []
    for block in section.blocks:
        if block.kind == "table":
            rows = block.value
            if len(rows) < 2:
                return []
            points = []
            for row in rows[1: 1 + limit]:
                if len(row) >= 3:
                    points.append(f"{row[0]}：{row[1]}，角色是{row[2]}")
            return points
    return []


def collect_rule_points(section: SectionNode | None) -> list[str]:
    if not section:
        return []
    mapping = []
    for child in section.children:
        paragraph = first_paragraph(child)
        if paragraph:
            mapping.append(f"{child.title.replace('规则 ', '')}：{paragraph}")
    return mapping[:5]


def set_text_frame_font(text_frame, size: int, color: tuple[int, int, int] = TITLE_COLOR, bold: bool = False) -> None:
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = FONT_NAME
            run.font.size = PptPt(size)
            run.font.bold = bold
            run.font.color.rgb = PptRGBColor(*color)


def add_background(slide, dark: bool = False) -> None:
    fill = DARK_BG if dark else LIGHT_BG
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PptRGBColor(*fill)
    bg.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, PptInches(0.2))
    accent.fill.solid()
    accent.fill.fore_color.rgb = PptRGBColor(*GOLD)
    accent.line.fill.background()


def add_title_slide(slide, title: str, subtitle: str) -> None:
    add_background(slide, dark=True)

    title_box = slide.shapes.add_textbox(PptInches(0.75), PptInches(1.2), PptInches(8.5), PptInches(1.8))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    set_text_frame_font(tf, size=28, color=WHITE, bold=True)

    subtitle_box = slide.shapes.add_textbox(PptInches(0.8), PptInches(3.1), PptInches(7.8), PptInches(1.3))
    tf2 = subtitle_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    set_text_frame_font(tf2, size=16, color=(224, 229, 236))

    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(0.8), PptInches(5.5), PptInches(2.3), PptInches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = PptRGBColor(*GOLD)
    badge.line.fill.background()
    badge_tf = badge.text_frame
    badge_tf.paragraphs[0].text = "内部使用 / 汇报版"
    badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    set_text_frame_font(badge_tf, size=12, color=DARK_BG, bold=True)


def add_bullet_slide(slide, title: str, bullets: list[str], subtitle: str | None = None) -> None:
    add_background(slide, dark=False)

    title_box = slide.shapes.add_textbox(PptInches(0.7), PptInches(0.6), PptInches(9.2), PptInches(0.6))
    title_tf = title_box.text_frame
    title_tf.paragraphs[0].text = title
    set_text_frame_font(title_tf, size=24, color=TITLE_COLOR, bold=True)

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(PptInches(0.75), PptInches(1.2), PptInches(5.5), PptInches(0.4))
        subtitle_tf = subtitle_box.text_frame
        subtitle_tf.paragraphs[0].text = subtitle
        set_text_frame_font(subtitle_tf, size=11, color=(154, 122, 28), bold=True)

    body = slide.shapes.add_textbox(PptInches(0.9), PptInches(1.8), PptInches(11.0), PptInches(5.0))
    tf = body.text_frame
    tf.word_wrap = True
    tf.margin_left = PptInches(0.03)
    tf.margin_right = PptInches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.clear()

    for idx, bullet in enumerate(bullets):
        paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.bullet = True
        paragraph.space_after = PptPt(8)

    set_text_frame_font(tf, size=18, color=TITLE_COLOR)

    note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(9.3), PptInches(0.55), PptInches(2.2), PptInches(0.5))
    note.fill.solid()
    note.fill.fore_color.rgb = PptRGBColor(240, 235, 216)
    note.line.fill.background()
    note_tf = note.text_frame
    note_tf.paragraphs[0].text = "当前版本"
    note_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    set_text_frame_font(note_tf, size=11, color=(154, 122, 28), bold=True)


def export_ppt(root: SectionNode) -> None:
    slide_plan = build_slide_plan(root)

    global prs
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)

    blank_layout = prs.slide_layouts[6]

    for slide_data in slide_plan:
        slide = prs.slides.add_slide(blank_layout)
        if slide_data.get("type") == "title":
            add_title_slide(slide, str(slide_data["title"]), str(slide_data["subtitle"]))
        else:
            add_bullet_slide(
                slide,
                title=str(slide_data["title"]),
                bullets=[str(item) for item in slide_data.get("bullets", [])],
                subtitle=str(slide_data["subtitle"]) if slide_data.get("subtitle") else None,
            )

    prs.save(PPT_OUTPUT)


def main() -> None:
    markdown_text = SOURCE_MD.read_text(encoding="utf-8")
    root = parse_markdown_sections(markdown_text)
    export_word(root)
    export_ppt(root)
    print(f"Generated: {WORD_OUTPUT}")
    print(f"Generated: {PPT_OUTPUT}")


if __name__ == "__main__":
    main()
