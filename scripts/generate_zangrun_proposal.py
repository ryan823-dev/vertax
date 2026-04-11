from __future__ import annotations

from datetime import date
from pathlib import Path

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.style import WD_STYLE_TYPE
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches as DocInches
from docx.shared import Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt as PptPt


OUTPUT_DIR = Path(r"D:\vertax\docs\proposals")
TODAY = date.today().isoformat()
PPT_PATH = OUTPUT_DIR / f"vertax-zangrun-proposal-{TODAY}.pptx"
DOCX_PATH = OUTPUT_DIR / f"vertax-zangrun-proposal-{TODAY}.docx"


BRAND_NAVY = RGBColor(0x0B, 0x1B, 0x2B)
BRAND_GOLD = RGBColor(0xD4, 0xAF, 0x37)
BRAND_IVORY = RGBColor(0xFF, 0xFC, 0xF7)
BRAND_TEXT = RGBColor(0x2A, 0x2A, 0x2A)
BRAND_MUTED = RGBColor(0x6B, 0x72, 0x80)
BRAND_LINE = RGBColor(0xE8, 0xE0, 0xD0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


TITLE = "把一家想出海的企业，变成一家能出海的企业"
SUBTITLE = "VertaX 赋能山东藏润环保科技有限公司出海整体方案"
TAGLINE = "专业出海官网（含 AI 员工） + VertaX AI 增长系统 + 官网/API 联动闭环"


def ensure_dir() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


def set_run_font(run, size: int, bold: bool = False, color: RGBColor = BRAND_TEXT) -> None:
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Microsoft YaHei"


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False,
                color: RGBColor = BRAND_TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_font(run, font_size, bold=bold, color=color)
    return box


def add_bullets(box, items, level_font_sizes=None, color: RGBColor = BRAND_TEXT):
    level_font_sizes = level_font_sizes or {0: 20, 1: 14}
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        if isinstance(item, tuple):
            level, text = item
        else:
            level, text = 0, item
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = level
        p.text = text
        p.alignment = PP_ALIGN.LEFT
        p.space_after = PptPt(6)
        for run in p.runs:
            set_run_font(run, level_font_sizes.get(level, 14), bold=(level == 0), color=color)
    return box


def add_title(slide, title: str, subtitle: str | None = None):
    add_textbox(slide, Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.7), title, font_size=24, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.62), Inches(1.02), Inches(11.8), Inches(0.4), subtitle, font_size=11, color=BRAND_MUTED)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(1.36), Inches(1.3), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = BRAND_GOLD
    line.line.fill.background()


def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_card(slide, left, top, width, height, title, body, number=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BRAND_LINE
    if number:
        badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(0.2), top + Inches(0.18), Inches(0.42), Inches(0.42))
        badge.fill.solid()
        badge.fill.fore_color.rgb = BRAND_GOLD
        badge.line.fill.background()
        add_textbox(slide, left + Inches(0.2), top + Inches(0.2), Inches(0.42), Inches(0.3), str(number), font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        title_left = left + Inches(0.75)
        title_width = width - Inches(0.95)
    else:
        title_left = left + Inches(0.25)
        title_width = width - Inches(0.5)
    add_textbox(slide, title_left, top + Inches(0.18), title_width, Inches(0.38), title, font_size=16, bold=True)
    add_textbox(slide, left + Inches(0.25), top + Inches(0.72), width - Inches(0.5), height - Inches(0.9), body, font_size=11, color=BRAND_MUTED)


def build_ppt() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_NAVY)
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.7), Inches(0.95), Inches(0.15), Inches(4.9))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BRAND_GOLD
    accent.line.fill.background()
    add_textbox(slide, Inches(1.15), Inches(1.05), Inches(11.1), Inches(1.5), TITLE, font_size=28, bold=True, color=WHITE)
    add_textbox(slide, Inches(1.15), Inches(2.7), Inches(10.8), Inches(0.8), SUBTITLE, font_size=18, color=WHITE)
    add_textbox(slide, Inches(1.15), Inches(3.55), Inches(10.8), Inches(0.8), TAGLINE, font_size=14, color=BRAND_GOLD)
    add_textbox(slide, Inches(1.15), Inches(5.75), Inches(5.0), Inches(0.35), "客户：山东藏润环保科技有限公司", font_size=12, color=WHITE)
    add_textbox(slide, Inches(1.15), Inches(6.1), Inches(5.0), Inches(0.35), f"日期：{TODAY}", font_size=12, color=WHITE)

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_IVORY)
    add_title(slide, "项目判断", "为什么藏润现在适合启动出海")
    add_textbox(
        slide,
        Inches(0.75),
        Inches(1.7),
        Inches(11.7),
        Inches(1.0),
        "藏润环保已经具备启动出海的产品与市场基础，但尚未形成真正可持续运行的出海体系。",
        font_size=22,
        bold=True,
    )
    add_card(slide, Inches(0.75), Inches(2.75), Inches(3.8), Inches(2.7), "产品具有核心竞争力", "核孔膜自发式气调保鲜箱具备明确技术壁垒，拥有专利、认证、标准参与和市场应用验证。")
    add_card(slide, Inches(4.75), Inches(2.75), Inches(3.8), Inches(2.7), "海外存在潜在市场需求", "果蔬、水产、冷链、绿色包装和降本提效场景具备全球共性需求，具备切入国际 B2B 市场的条件。")
    add_card(slide, Inches(8.75), Inches(2.75), Inches(3.8), Inches(2.7), "当前缺的是系统化出海能力", "企业并不缺产品，而是缺少路径判断、专业表达、官网承接、获客启动与组织沉淀。")

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_IVORY)
    add_title(slide, "内部准备不足", "企业还没真正具备出海起步能力")
    internal_items = [
        ("1. 路径不清", "不知道先做哪个市场、哪个渠道、哪种出海模式，方向不明。"),
        ("2. 资料不全", "企业介绍、产品资料、案例内容、技术说明准备不足，难以支撑海外启动。"),
        ("3. 表达不专业", "企业价值、产品优势和服务能力说不清，难以建立第一轮信任。"),
        ("4. 配套缺失", "认证、资质、进出口、关务、海外仓、收款等配套经验不足。"),
    ]
    positions = [(Inches(0.75), Inches(1.7)), (Inches(6.7), Inches(1.7)), (Inches(0.75), Inches(4.05)), (Inches(6.7), Inches(4.05))]
    for idx, ((title, body), (left, top)) in enumerate(zip(internal_items, positions), start=1):
        add_card(slide, left, top, Inches(5.8), Inches(1.9), title, body, number=idx)

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_IVORY)
    add_title(slide, "外部启动困难", "市场还没有真正跑起来")
    external_items = [
        ("5. 获客无从下手", "SEO、广告、社媒、展会、平台、主动开发之间缺乏清晰切入顺序。"),
        ("6. 投入门槛高", "建团队、跑展会、做推广都要持续投入，前期成本和试错压力大。"),
        ("7. 冷启动困难", "没有稳定渠道、品牌基础和海外客户触达能力，第一批客户难启动。"),
        ("8. 人员依赖", "过度依赖业务员个人能力，一旦换人，客户推进和资源积累容易中断。"),
    ]
    for idx, ((title, body), (left, top)) in enumerate(zip(external_items, positions), start=5):
        add_card(slide, left, top, Inches(5.8), Inches(1.9), title, body, number=idx)

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_IVORY)
    add_title(slide, "VertaX 提供的三种能力", "我们不是单卖模块，而是在帮助企业获得真正的出海能力")
    capability_cards = [
        ("专业化表达", "让专业，被客户一眼看见\n把企业、产品与服务能力，转化为清晰、专业、可信的海外表达。"),
        ("体系化启动", "让出海营销与获客真正启动\n围绕内容、渠道与客户触达，推动出海动作从想法走向执行。"),
        ("组织化沉淀", "让每一次动作都成为能力\n把资料、内容、线索与流程沉淀为可复用的组织资产。"),
    ]
    for i, (title, body) in enumerate(capability_cards):
        add_card(slide, Inches(0.75 + i * 4.15), Inches(2.1), Inches(3.9), Inches(3.2), title, body)

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_IVORY)
    add_title(slide, "整体方案总览", "藏润环保出海方案由三层组成")
    layers = [
        ("第一层：专业出海官网", "面向海外客户的品牌与获客前台，承载专业表达、产品展示、案例建立与询盘承接。"),
        ("第二层：VertaX AI 增长系统", "负责知识沉淀、内容生成、客户发现、主动触达、协同推进与结果判断。"),
        ("第三层：官网与系统 API 打通", "让内容、线索、页面更新与结果反馈形成闭环，官网不再只是静态展示页。"),
    ]
    top = Inches(1.8)
    for title, body in layers:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.1), top, Inches(11.1), Inches(1.3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = BRAND_LINE
        add_textbox(slide, Inches(1.4), top + Inches(0.15), Inches(4.0), Inches(0.35), title, font_size=16, bold=True)
        add_textbox(slide, Inches(4.7), top + Inches(0.15), Inches(7.0), Inches(0.65), body, font_size=12, color=BRAND_MUTED)
        top += Inches(1.55)

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_IVORY)
    add_title(slide, "带 AI 员工的专业出海官网", "这不是普通展示站，而是海外业务前台")
    left_box = slide.shapes.add_textbox(Inches(0.85), Inches(1.75), Inches(5.6), Inches(4.8))
    add_bullets(left_box, [
        "多语种官网框架，适配海外客户浏览习惯",
        "围绕企业、产品、认证、案例和应用场景重建表达",
        "产品页与解决方案页按客户痛点和合作模式组织",
        "AI 员工 24 小时在线，负责问答、筛选和线索收集",
        "天然支持 SEO/AEO/GEO 内容布局，具备长期自然获客基础",
    ], level_font_sizes={0: 18})
    right = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.9), Inches(5.4), Inches(4.4))
    right.fill.solid()
    right.fill.fore_color.rgb = BRAND_NAVY
    right.line.fill.background()
    add_textbox(slide, Inches(7.35), Inches(2.25), Inches(4.6), Inches(0.45), "官网前台角色", font_size=16, bold=True, color=BRAND_GOLD)
    add_textbox(slide, Inches(7.35), Inches(3.0), Inches(4.6), Inches(2.8),
                "1. 展示专业能力\n2. 承接海外访问\n3. 由 AI 员工完成首轮问答\n4. 自动采集客户需求与联系方式\n5. 将线索回流到 VertaX 系统继续推进",
                font_size=14, color=WHITE)

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_IVORY)
    add_title(slide, "VertaX 六大模块", "六大模块分别服务于企业的三种出海能力")
    rows = [
        ("专业化表达", "知识引擎 + 声量枢纽", "沉淀企业知识、统一专业表达，并持续输出市场声量与信任感。"),
        ("体系化启动", "增长系统 + 获客雷达", "持续做内容、布局 SEO/AEO/GEO，并主动识别客户、发现线索、推动触达。"),
        ("组织化沉淀", "推进中台 + 决策中心", "承接协同互动与内容校正，沉淀过程数据并辅助经营判断。"),
    ]
    y = Inches(1.9)
    for ability, modules, body in rows:
        add_card(slide, Inches(0.9), y, Inches(2.2), Inches(1.35), ability, "")
        add_card(slide, Inches(3.35), y, Inches(2.65), Inches(1.35), modules, "")
        add_card(slide, Inches(6.25), y, Inches(6.0), Inches(1.35), "模块价值", body)
        y += Inches(1.65)

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_IVORY)
    add_title(slide, "官网与系统 API 联动", "让内容、线索和结果形成真正的闭环")
    api_steps = [
        ("内容下发", "VertaX 生成的行业文章、FAQ、产品内容、案例内容，可通过 API 推送到官网对应栏目。"),
        ("线索回流", "官网表单、AI 员工对话、询盘提交的数据，自动回流到 VertaX 系统进入线索池。"),
        ("页面持续优化", "随着知识库和内容系统不断产出，官网页面持续同步更新，而不是一次性建完结束。"),
        ("结果反馈", "官网访问、内容表现、询盘质量和客户反馈回流系统，用于持续判断和优化。"),
    ]
    for idx, (title, body) in enumerate(api_steps, start=1):
        add_card(slide, Inches(0.9 + ((idx - 1) % 2) * 6.1), Inches(1.85 + ((idx - 1) // 2) * 2.15), Inches(5.35), Inches(1.75), title, body, number=idx)

    # Slide 10
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_IVORY)
    add_title(slide, "为什么这套方案适合藏润环保", "产品力已具备，当前最需要的是系统化启动能力")
    left = slide.shapes.add_textbox(Inches(0.85), Inches(1.75), Inches(5.9), Inches(4.8))
    add_bullets(left, [
        "2017 年成立，总部烟台，具备清华合作技术背景",
        "核孔膜自发式气调保鲜箱，具备明确差异化技术路线",
        "拥有 54 项专利，其中包含发明专利",
        "参与《限制快递过度包装要求》国家标准相关工作",
        "已获得三体系、绿色产品、碳足迹、低碳及欧盟 SGS 等认证",
        "已在果蔬、鲜活水产等场景应用，并与顺丰、中国邮政、京东等形成合作和验证",
    ], level_font_sizes={0: 16})
    right = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.1), Inches(1.9), Inches(5.1), Inches(4.4))
    right.fill.solid()
    right.fill.fore_color.rgb = WHITE
    right.line.color.rgb = BRAND_LINE
    add_textbox(slide, Inches(7.45), Inches(2.25), Inches(4.2), Inches(0.4), "建议优先切入方向", font_size=16, bold=True)
    add_textbox(slide, Inches(7.45), Inches(2.95), Inches(4.2), Inches(2.6),
                "• 以 B2B 为主，不先走 B2C\n• 优先寻找海外代理商、渠道商、进口商\n• 突出绿色包装、保鲜能力、无冰无冷链、降本提效\n• 验证“核心材料/核心能力输出 + 本地合作交付”路径",
                font_size=14, color=BRAND_MUTED)

    # Slide 11
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_IVORY)
    add_title(slide, "实施路径", "分阶段推进，从搭建到启动，再到沉淀")
    phases = [
        ("阶段一\n基础搭建期", "1-4 周", "梳理企业资料，搭建知识库，明确目标市场与客户画像，完成官网结构规划与核心表达体系。"),
        ("阶段二\n启动运行期", "5-8 周", "官网上线，AI 员工上线，VertaX 系统部署完成，API 打通完成，启动第一批内容与客户发现。"),
        ("阶段三\n持续优化期", "9-12 周", "持续更新内容、启动定向触达、跟踪询盘结果，完善知识库、线索流程与经营复盘。"),
    ]
    for i, (name, period, body) in enumerate(phases):
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75 + i * 4.2), Inches(2.0), Inches(3.85), Inches(3.6))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = BRAND_LINE
        add_textbox(slide, Inches(1.0 + i * 4.2), Inches(2.25), Inches(3.3), Inches(0.9), name, font_size=18, bold=True)
        add_textbox(slide, Inches(1.0 + i * 4.2), Inches(3.0), Inches(2.2), Inches(0.35), period, font_size=12, bold=True, color=BRAND_GOLD)
        add_textbox(slide, Inches(1.0 + i * 4.2), Inches(3.45), Inches(3.2), Inches(1.65), body, font_size=11, color=BRAND_MUTED)

    # Slide 12
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_IVORY)
    add_title(slide, "项目最终交付", "交付的不只是一个网站，而是一套能持续运转的出海系统")
    deliverables = [
        "一套带 AI 员工的专业出海官网",
        "一套部署完成的 VertaX AI 增长系统",
        "一套官网与系统之间的 API 联动机制",
        "一套企业专属知识与表达底座",
        "一套持续运行的内容、获客与线索闭环机制",
    ]
    box = slide.shapes.add_textbox(Inches(0.95), Inches(1.8), Inches(11.3), Inches(3.8))
    add_bullets(box, deliverables, level_font_sizes={0: 20})
    ribbon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(5.85), Inches(11.2), Inches(0.7))
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = BRAND_NAVY
    ribbon.line.fill.background()
    add_textbox(slide, Inches(1.2), Inches(6.0), Inches(10.5), Inches(0.3), "官网负责对外表达，系统负责持续运转，API 负责让两者真正连起来。", font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Slide 13
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BRAND_NAVY)
    add_textbox(slide, Inches(0.85), Inches(0.7), Inches(5.8), Inches(0.6), "年度合作报价", font_size=28, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.85), Inches(1.35), Inches(6.6), Inches(0.4), "专业出海官网 + 带 AI 员工的前台 + VertaX 系统 + API 打通 + 年度运营支撑", font_size=13, color=WHITE)
    price_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(2.0), Inches(4.6), Inches(2.6))
    price_box.fill.solid()
    price_box.fill.fore_color.rgb = BRAND_GOLD
    price_box.line.fill.background()
    add_textbox(slide, Inches(1.2), Inches(2.45), Inches(3.8), Inches(0.45), "人民币", font_size=16, bold=True, color=BRAND_NAVY, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.1), Inches(2.9), Inches(4.0), Inches(0.9), "198,000", font_size=32, bold=True, color=BRAND_NAVY, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.2), Inches(3.85), Inches(3.8), Inches(0.35), "元 / 年", font_size=16, bold=True, color=BRAND_NAVY, align=PP_ALIGN.CENTER)
    right_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.0), Inches(2.0), Inches(6.0), Inches(3.3))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = WHITE
    right_box.line.fill.background()
    add_textbox(slide, Inches(6.35), Inches(2.3), Inches(5.2), Inches(0.35), "报价包含", font_size=16, bold=True)
    box = slide.shapes.add_textbox(Inches(6.35), Inches(2.8), Inches(5.1), Inches(2.0))
    add_bullets(box, [
        "专业出海官网搭建",
        "带 AI 员工的智能接待能力部署",
        "VertaX 系统部署与初始化配置",
        "官网与系统 API 打通",
        "企业知识底座建设",
        "年度运营支撑与优化建议",
    ], level_font_sizes={0: 15})
    add_textbox(slide, Inches(6.35), Inches(5.15), Inches(5.1), Inches(0.55), "说明：不含大规模广告投放费用、第三方平台采购费及新增私有化硬件成本。", font_size=10, color=BRAND_MUTED)

    prs.save(str(PPT_PATH))


def set_doc_run_font(run, size: int, bold: bool = False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = "Microsoft YaHei"
    run._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")


def add_doc_paragraph(doc: Document, text: str, size=11, bold=False, style=None):
    p = doc.add_paragraph(style=style)
    run = p.add_run(text)
    set_doc_run_font(run, size=size, bold=bold)
    return p


def add_doc_bullets(doc: Document, items: list[str]) -> None:
    for item in items:
        p = doc.add_paragraph(style="List Bullet")
        run = p.add_run(item)
        set_doc_run_font(run, 10)


def build_docx() -> None:
    doc = Document()
    section = doc.sections[0]
    section.top_margin = DocInches(0.7)
    section.bottom_margin = DocInches(0.7)
    section.left_margin = DocInches(0.85)
    section.right_margin = DocInches(0.85)

    styles = doc.styles
    if "BodyCN" not in styles:
        style = styles.add_style("BodyCN", WD_STYLE_TYPE.PARAGRAPH)
        style.font.name = "Microsoft YaHei"
        style._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
        style.font.size = Pt(10.5)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run(TITLE)
    set_doc_run_font(run, 22, bold=True)

    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = subtitle.add_run(SUBTITLE)
    set_doc_run_font(run, 13)

    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = meta.add_run(f"{TAGLINE}\n日期：{TODAY}")
    set_doc_run_font(run, 10)

    doc.add_paragraph("")

    sections = [
        ("一、项目判断", [
            "藏润环保已经具备启动出海的产品与市场基础，但尚未形成真正可持续运行的出海体系。",
            "企业当前并不缺产品，而是缺少系统化的路径判断、专业表达、官网承接、内容启动、客户发现与组织沉淀能力。",
            "VertaX 本次提供的不是单独的网站或软件，而是一套由专业出海官网、VertaX AI 增长系统以及 API 联动组成的整体方案。"
        ]),
        ("二、藏润启动出海的基础条件", [
            "产品具有核心竞争力：核孔膜自发式气调保鲜箱具备明确技术壁垒与差异化价值。",
            "企业有硬实力支撑：与清华大学合作，拥有 54 项专利，并参与国家标准相关工作。",
            "行业与市场验证充分：在果蔬、鲜活水产、绿色包装场景已形成应用，并与顺丰、中国邮政、京东等形成合作验证。",
            "当前最需要的是把已有优势转化为海外市场可理解、可相信、可合作的能力体系。"
        ]),
        ("三、企业当前面临的两大部分、八大挑战", [
            "第一部分：内部准备不足。路径不清、资料不全、表达不专业、配套缺失。",
            "第二部分：外部启动困难。获客无从下手、投入门槛高、冷启动困难、人员依赖严重。",
            "这些问题不是靠单独招聘或一次性投入就能解决，而需要系统化机制持续支撑。"
        ]),
        ("四、VertaX 帮企业建立的三种能力", [
            "1. 专业化表达：让专业，被客户一眼看见。",
            "2. 体系化启动：让出海营销与获客真正启动。",
            "3. 组织化沉淀：让每一次动作都成为能力。"
        ]),
        ("五、整体方案结构", [
            "第一层：专业出海官网。作为企业面向海外市场的品牌与获客前台，负责专业表达、产品展示、案例建立、询盘承接和信任建立。",
            "第二层：VertaX AI 增长系统。负责知识沉淀、内容生成、客户发现、主动触达、协同推进与结果判断。",
            "第三层：官网与系统 API 打通。实现内容下发、线索回流、页面更新与数据反馈，形成闭环。"
        ]),
        ("六、带 AI 员工的专业出海官网", [
            "官网将不再只是展示页，而是面向海外客户的业务前台。",
            "官网支持多语种框架，围绕企业、产品、认证、案例、行业场景和合作模式建立表达结构。",
            "官网内嵌 AI 员工，可 24 小时完成客户问答、需求引导、联系方式采集和首轮线索筛选。",
            "官网天然支持后续 SEO、AEO、GEO 布局，能够持续承接内容和自然流量。"
        ]),
        ("七、VertaX 六大模块如何服务藏润环保", [
            "知识引擎 + 声量枢纽：服务于专业化表达。",
            "增长系统 + 获客雷达：服务于体系化启动。",
            "推进中台 + 决策中心：服务于组织化沉淀。",
            "六大模块并不是孤立功能，而是共同支撑官网、内容、线索、协同与经营判断的底层体系。"
        ]),
        ("八、官网与系统 API 打通的价值", [
            "VertaX 生成的文章、FAQ、产品内容和案例内容可同步到官网。",
            "官网表单、AI 员工对话和询盘线索可自动回流到系统。",
            "官网页面可随着知识库和内容系统持续更新，而不是一次性交付。",
            "官网表现、内容结果和询盘反馈可回流到系统，辅助下一轮优化。"
        ]),
        ("九、针对藏润环保的专项落地建议", [
            "优先以 B2B 为主，不建议当前阶段优先做 B2C。",
            "优先面向海外代理商、渠道商、进口商及冷链、生鲜、包装相关 B2B 合作方。",
            "重点表达绿色包装、保鲜能力、无冰无冷链、运输成本优化和标准化能力。",
            "逐步验证“核心能力输出 + 海外本地合作交付”的业务路径。"
        ]),
        ("十、实施节奏", [
            "第一阶段（1-4 周）：梳理资料、搭建知识库、完成官网结构规划和核心表达体系。",
            "第二阶段（5-8 周）：官网上线、AI 员工上线、VertaX 系统部署完成、API 打通完成、启动首轮内容和线索。",
            "第三阶段（9-12 周）：持续更新内容、启动主动获客与定向触达、跟踪结果并沉淀流程。"
        ]),
        ("十一、项目最终交付", [
            "一套带 AI 员工的专业出海官网。",
            "一套部署完成的 VertaX AI 增长系统。",
            "一套官网与系统之间的 API 联动机制。",
            "一套企业专属知识与表达底座。",
            "一套持续运行的内容、获客与线索闭环机制。"
        ]),
    ]

    for heading, paragraphs in sections:
        doc.add_heading(heading, level=1)
        for paragraph in paragraphs:
            add_doc_paragraph(doc, paragraph, style="BodyCN")

    doc.add_heading("十二、合作方式与报价", level=1)
    add_doc_paragraph(
        doc,
        "本项目采用年度合作方式，由 VertaX 提供“专业出海官网 + 带 AI 员工的海外接待前台 + VertaX AI 增长系统 + API 打通 + 持续运营支撑”的整体解决方案。",
        style="BodyCN",
    )

    quote = doc.add_paragraph()
    quote.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = quote.add_run("报价：人民币 198,000 元 / 年")
    set_doc_run_font(run, 18, bold=True)

    table = doc.add_table(rows=1, cols=2)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    hdr = table.rows[0].cells
    hdr[0].text = "报价包含"
    hdr[1].text = "说明"
    items = [
        ("专业出海官网搭建", "含官网结构规划、核心页面搭建、多语种基础框架"),
        ("带 AI 员工的智能接待能力部署", "含问答、需求引导、联系方式采集与首轮筛选"),
        ("VertaX 系统部署与初始化配置", "含知识引擎、增长系统、获客雷达、推进中台、决策中心等"),
        ("官网与系统 API 打通", "含内容同步、线索回流与数据联动"),
        ("企业知识底座建设", "含企业、产品、资质、案例、FAQ 等结构化沉淀"),
        ("年度运营支撑", "含内容推进、阶段复盘与优化建议"),
    ]
    for left, right in items:
        row = table.add_row().cells
        row[0].text = left
        row[1].text = right

    doc.add_paragraph("")
    add_doc_paragraph(
        doc,
        "以上报价不含大规模广告投放费用、第三方平台会员采购费用及企业新增私有化硬件服务器成本；如需私有化部署或特殊集成，可另行评估。",
        style="BodyCN",
    )

    doc.add_heading("十三、客户需配合提供的基础资料", level=1)
    add_doc_bullets(doc, [
        "企业介绍、品牌资料、核心产品资料",
        "认证证书、检测报告、专利与标准参与信息",
        "典型案例、合作客户、应用场景资料",
        "现有英文资料、图片、视频、PPT 或 PDF",
        "优先市场、优先语言、优先产品和目标客户类型",
    ])

    doc.add_paragraph("")
    end = doc.add_paragraph()
    end.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = end.add_run("VertaX 不是另一个单点工具，而是帮助企业从“想出海”进入“能出海、会出海、持续出海”的系统能力。")
    set_doc_run_font(run, 12, bold=True)

    doc.save(str(DOCX_PATH))


def main() -> None:
    ensure_dir()
    build_ppt()
    build_docx()
    print(PPT_PATH)
    print(DOCX_PATH)


if __name__ == "__main__":
    main()
